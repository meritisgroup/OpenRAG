from .tools import list_folders

from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from ..utils.splitter import MarkdownHeaderTextSplitter
import fitz
import os


def get_position(shape):
    return shape.left + shape.top


def open_ppt(path_ppt: str) -> str:
    """
    Permet d'extraire tout le texte d'un ppt enregistré au chemin path_ppt
    """
    presentation = Presentation(path_ppt)

    content = ""

    for slide_num, slide in enumerate(presentation.slides, start=1):

        textboxes = [
            shape
            for shape in slide.shapes
            if hasattr(shape, "text_frame") and shape.text_frame is not None
        ]

        sorted_textboxes = sorted(textboxes, key=get_position)

        for shape in sorted_textboxes:
            textbox_text = shape.text_frame.text

            content += f"{textbox_text}\n\n"

    return content


def open_pdf(path_pdf: str) -> str:
    """
    Permet d'extraire tout le texte d'un PDF enregistré au chemin path_pdf
    """

    with fitz.open(path_pdf) as document:
       return "".join(page.get_text("text") for page in document)


def open_excel(path_excel: str) -> str:
    """
    Permet d'extraire tout le texte d'un fichier Excel (.xlsx) enregistré au chemin path_xlsx.
    """
    workbook = load_workbook(filename=path_excel, data_only=True)
    content = ""

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            content += " ".join(str(cell) for cell in row if cell is not None) + "\n"

    return content


def open_word(path_word: str) -> str:
    """
    Permet d'extraire tout le texte d'un fichier Word (.docx) enregistré au chemin path_docx.
    """
    document = Document(path_word)
    content = ""

    for paragraph in document.paragraphs:
        content += paragraph.text + "\n"

    return content


def open_txt(path_txt: str) -> str:
    """
    Permet d'extraire tout le texte d'un fichier .txt enregistré au chemin path_txt
    """
    with open(path_txt, "r", encoding="utf-8") as f:
        content = f.read()

    return content


def open_doc_with_save(path_file: str, path_to_save: str, overwrite=False) -> str:
    """
    Permet de récupérer les informations textuels d'un PDF, PPT, Word, Excel ou Txt.
    Si le paramètre save est set sur True alors on enregistre le résultat sous un format txt au chemin path_to_save
    """
    if path_to_save[-1] != "/":
        path_to_save += "/"

    if ".xlsx" in path_file:
        content = open_excel(path_file)

    elif ".docx" in path_file:
        content = open_word(path_file)

    elif ".pptx" in path_file:
        content = open_ppt(path_file)

    elif ".pdf" in path_file:
        content = open_pdf(path_file)

    elif ".txt" in path_file:
        content = open_txt(path_file)

    else:
        content = "\nExtension de fichier inconnu\n"

    name_doc = (path_file.split("/"))[-1]
    name_doc_sans_extension = (name_doc.split("."))[0]
    final_path = path_to_save + name_doc_sans_extension + ".txt"

    if overwrite or final_path not in list_folders(path_to_save):
        with open(final_path, "w", encoding="utf-8") as f:
            f.write(content)

    return content


def open_doc_without_save(path_file: str) -> str:
    """
    Permet de récupérer les informations textuels d'un PDF, PPT, Word, Excel ou Txt.
    Si le paramètre save est set sur True alors on enregistre le résultat sous un format txt au chemin path_to_save
    """

    if ".xlsx" in path_file:
        content = open_excel(path_file)

    elif ".docx" in path_file:
        content = open_word(path_file)

    elif ".pptx" in path_file:
        content = open_ppt(path_file)

    elif ".pdf" in path_file:
        content = open_pdf(path_file)

    elif ".txt" in path_file or ".md" in path_file:
        content = open_txt(path_file)
        
    else:
        content = "\nExtension de fichier inconnu\n"

    return content


class Opener:
    def __init__(self, save=False, overwrite=False) -> None:
        self._save = save
        self.overwrite = overwrite

    def open_doc(self, path_file, path_to_save=""):
        return open_doc_without_save(path_file)

    @property
    def save(self):
        print(f'The opener is set on "{self._save}" mode')
        return self._save

    @save.setter
    def save(self, value) -> None:
        self._save = value
